"""Build GrantGenie.pptx from docs/client/GrantGenie-Presentation-Outline.md.

Uses python-pptx to produce a 16:9 deck matching the design tokens defined in
the outline: indigo primary, coral accent, warm white background, Inter-style
typeface. Falls back to Calibri (always available) when custom fonts are
absent so the deck opens cleanly on any machine.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------

INDIGO = RGBColor(0x3D, 0x2E, 0xA0)
INDIGO_LIGHT = RGBColor(0x6E, 0x60, 0xD5)
CORAL = RGBColor(0xFF, 0x6B, 0x5C)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
WARM_WHITE = RGBColor(0xFA, 0xFA, 0xF7)
NEAR_BLACK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def add_background(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 18,
    color: RGBColor = NEAR_BLACK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font: str = BODY_FONT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    *,
    size: int = 18,
    color: RGBColor = NEAR_BLACK,
    bullet_color: RGBColor = CORAL,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        bullet = p.add_run()
        bullet.text = "●  "
        bullet.font.name = BODY_FONT
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color
        bullet.font.bold = True
        run = p.add_run()
        run.text = item
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    line: RGBColor | None = None,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False


def add_rounded(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: RGBColor,
    line: RGBColor | None = None,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.adjustments[0] = 0.12
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False


def add_footer(slide, page: int, total: int) -> None:
    add_text_box(
        slide,
        0.5,
        7.05,
        12.3,
        0.3,
        f"GrantGenie · v1.0 · July 2026            {page} / {total}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.LEFT,
    )


def add_title(slide, text: str, *, color: RGBColor = NEAR_BLACK, top: float = 0.55) -> None:
    add_text_box(
        slide,
        0.7,
        top,
        12.0,
        1.0,
        text,
        size=36,
        color=color,
        bold=True,
        font=TITLE_FONT,
    )


def add_subtitle(slide, text: str, *, color: RGBColor = MUTED, top: float = 1.5) -> None:
    add_text_box(
        slide,
        0.7,
        top,
        12.0,
        0.5,
        text,
        size=18,
        color=color,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_01_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Indigo gradient background (simulated with two stacked rectangles)
    add_background(slide, INDIGO)
    add_rect(slide, 0, 0, 13.333, 7.5, INDIGO)
    add_rect(slide, 0, 0, 5.0, 7.5, INDIGO_LIGHT)

    # Coral accent line on the right
    add_rect(slide, 11.5, 0, 0.06, 7.5, CORAL)

    add_text_box(
        slide,
        0.8,
        2.4,
        11.5,
        1.6,
        "GrantGenie",
        size=84,
        color=WHITE,
        bold=True,
        font=TITLE_FONT,
    )
    add_text_box(
        slide,
        0.8,
        4.1,
        11.5,
        0.7,
        "AI-assisted grant discovery and drafting for small nonprofits",
        size=22,
        color=WHITE,
    )
    add_text_box(
        slide,
        0.8,
        6.4,
        11.5,
        0.4,
        "v1.0 · July 2026",
        size=12,
        color=WARM_WHITE,
    )


def slide_02_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "Small nonprofits lose funding to three predictable failure modes.")
    add_subtitle(slide, "Each is solvable — and each compounds when not addressed.")

    tiles = [
        ("1", "They don't know the grants exist.", "A typical state has 800+ active RFPs at any time."),
        ("2", "They spend 80% of proposal time on boilerplate.", "Re-typing the org's mission, programs, and metrics for every application."),
        ("3", "They submit unfocused proposals.", "Missed funder language, missed format, missed priorities."),
    ]
    for i, (num, headline, body) in enumerate(tiles):
        x = 0.7 + i * 4.1
        add_rounded(slide, x, 2.4, 3.9, 3.6, WHITE, line=MUTED)
        add_text_box(slide, x + 0.3, 2.6, 0.6, 0.6, num, size=44, color=CORAL, bold=True)
        add_text_box(slide, x + 0.3, 3.3, 3.3, 1.4, headline, size=18, color=NEAR_BLACK, bold=True)
        add_text_box(slide, x + 0.3, 4.6, 3.3, 1.4, body, size=14, color=MUTED)

    add_text_box(
        slide,
        0.7,
        6.3,
        12.0,
        0.5,
        "Industry benchmarks: 60–80 h/quarter on discovery · 12–25 h per 10-page draft · 15–25% pass-through rate",
        size=14,
        color=MUTED,
    )
    add_footer(slide, 2, 15)


def slide_03_why_now(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "The philanthropic sector is growing; staffing is not.")
    add_subtitle(slide, "AI + SaaS is now mature enough to deploy safely for regulated work.")

    kpis = [
        ("$500B+", "U.S. philanthropic disbursement per year"),
        ("+8–12% YoY", "RFP volume growth"),
        ("~10%", "Share flowing to small/medium nonprofits"),
    ]
    for i, (val, label) in enumerate(kpis):
        x = 0.7 + i * 4.1
        add_rounded(slide, x, 2.3, 3.9, 2.5, INDIGO)
        add_text_box(slide, x, 2.6, 3.9, 1.0, val, size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.3, 3.8, 3.3, 0.8, label, size=14, color=WARM_WHITE, align=PP_ALIGN.CENTER)

    add_bullets(
        slide,
        0.7,
        5.2,
        12.0,
        1.5,
        [
            "AI model quality, citation-grounded generation, and tenant-isolated SaaS are now mature enough for regulated work.",
            "Cloud + OpenTelemetry + eval-gated LLMs are the new normal.",
            "A new 'AI for good' SaaS category is opening up.",
        ],
        size=14,
    )
    add_footer(slide, 3, 15)


def slide_04_maya(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "Meet Maya — our primary user.")
    add_subtitle(slide, "A grant writer at a small nonprofit.")

    # Left persona card
    add_rounded(slide, 0.7, 2.3, 6.0, 4.2, WHITE, line=MUTED)
    add_text_box(slide, 1.0, 2.5, 5.4, 0.5, "Persona", size=12, color=CORAL, bold=True)
    add_text_box(
        slide,
        1.0,
        3.0,
        5.4,
        3.3,
        "Grant writer at a small nonprofit\n1–50 staff, $250k–$5M budget\n10–30 active applications per year\n30+ PDFs of institutional knowledge — no one reads them",
        size=16,
        color=NEAR_BLACK,
    )

    # Right day-in-life
    add_rounded(slide, 7.0, 2.3, 5.6, 4.2, WHITE, line=MUTED)
    add_text_box(slide, 7.3, 2.5, 5.0, 0.5, "A week in Maya's life", size=12, color=CORAL, bold=True)
    days = [("Mon", "Discovery"), ("Tue", "Drafting"), ("Wed", "Review with ED"), ("Thu", "Revisions"), ("Fri", "Submit + log")]
    for i, (day, task) in enumerate(days):
        y = 3.0 + i * 0.65
        add_text_box(slide, 7.3, y, 0.9, 0.5, day, size=18, color=INDIGO, bold=True)
        add_text_box(slide, 8.4, y, 4.0, 0.5, task, size=16, color=NEAR_BLACK)
    add_footer(slide, 4, 15)


def slide_05_solution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "One platform. Five steps. Cited at every step.")

    steps = ["Discover", "Match", "Draft", "Review", "Submit + Track"]
    node_w = 2.2
    gap = 0.18
    total_w = len(steps) * node_w + (len(steps) - 1) * gap
    start_x = (13.333 - total_w) / 2

    for i, label in enumerate(steps):
        x = start_x + i * (node_w + gap)
        add_rounded(slide, x, 3.4, node_w, 1.0, INDIGO)
        add_text_box(slide, x, 3.4, node_w, 1.0, label, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Arrows between nodes
    for i in range(len(steps) - 1):
        x = start_x + (i + 1) * node_w + i * gap + 0.01
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(3.75), Inches(gap), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = CORAL
        arrow.line.fill.background()

    add_text_box(
        slide,
        0.7,
        5.0,
        12.0,
        0.5,
        "Search only eligible grants. Get a full draft in 6 seconds with citations. Collaborate with reviewers. Track every deadline.",
        size=18,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 5, 15)


def slide_06_how_it_works(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "Behind the scenes — when you click 'Draft'.")

    # Left bullets
    add_bullets(
        slide,
        0.7,
        2.2,
        7.0,
        4.5,
        [
            "Your org profile + boilerplate + selected grant → AI Service",
            "Multi-model cost-aware router picks the cheapest adequate model",
            "RAG retrieves your top-5 relevant chunks; reranks to top-3",
            "LLM streams a 5-section draft with citations per claim",
            "Eval gates (relevance ≥ 0.85, faithfulness ≥ 0.90) block low-quality output",
        ],
        size=16,
    )

    # Right architecture sketch
    box_x = 8.4
    add_rounded(slide, box_x, 2.4, 4.2, 0.9, INDIGO)
    add_text_box(slide, box_x, 2.4, 4.2, 0.9, "User clicks Draft", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.2), Inches(3.35), Inches(0.4), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = CORAL
    arrow.line.fill.background()

    add_rounded(slide, box_x, 3.85, 4.2, 1.2, INDIGO_LIGHT)
    add_text_box(slide, box_x + 0.1, 3.95, 4.0, 1.0, "AI Service\n(multi-model router + eval gates)", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.2), Inches(5.10), Inches(0.4), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = CORAL
    arrow2.line.fill.background()

    add_rounded(slide, box_x, 5.55, 4.2, 1.0, CORAL)
    add_text_box(slide, box_x, 5.55, 4.2, 1.0, "Drafted proposal\nwith citations", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 6, 15)


def slide_07_trust(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "Built for the work that gets audited.")

    pillars = [
        ("Every claim cited", "Hover any AI sentence → see source document, page, and chunk."),
        ("Eval-gated quality", "Relevance ≥ 0.85 and faithfulness ≥ 0.90 required to surface."),
        ("Audit-ready logs", "Append-only audit log with before/after diff; triggers block any change."),
    ]
    for i, (head, body) in enumerate(pillars):
        x = 0.7 + i * 4.1
        add_rounded(slide, x, 2.4, 3.9, 3.6, WHITE, line=INDIGO)
        # Accent bar
        add_rect(slide, x, 2.4, 3.9, 0.1, INDIGO)
        add_text_box(slide, x + 0.3, 2.7, 3.3, 1.0, head, size=20, color=INDIGO, bold=True)
        add_text_box(slide, x + 0.3, 3.7, 3.3, 2.0, body, size=14, color=NEAR_BLACK)

    add_text_box(
        slide,
        0.7,
        6.3,
        12.0,
        0.5,
        "Multi-tenant isolation at the database layer makes cross-tenant data leaks structurally impossible.",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 7, 15)


def slide_08_ux(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "A working session in 60 seconds.")
    add_subtitle(slide, "Three moments Maya sees every week.")

    panels = [
        ("Discovery", "Only eligible grants shown.\n14 results in 3.2 s."),
        ("Drafting", "5-section draft, citations\ninline, eval scores visible."),
        ("Tracking", "Email + in-app reminders\nat 14, 7, 1 days."),
    ]
    for i, (head, body) in enumerate(panels):
        x = 0.7 + i * 4.1
        # Mock screen
        add_rounded(slide, x, 2.3, 3.9, 3.6, WHITE, line=MUTED)
        add_rect(slide, x + 0.2, 2.45, 3.5, 0.4, INDIGO)
        add_text_box(slide, x + 0.3, 2.45, 3.3, 0.4, "GrantGenie", size=12, color=WHITE, bold=True)
        # Mock content blocks
        for j in range(3):
            add_rect(slide, x + 0.3, 3.05 + j * 0.7, 3.4, 0.5, WARM_WHITE)
        add_rect(slide, x + 0.3, 5.4, 1.5, 0.3, CORAL)
        add_text_box(slide, x + 0.3, 6.05, 3.3, 0.5, head, size=18, color=INDIGO, bold=True)
        add_text_box(slide, x + 0.3, 6.5, 3.3, 0.5, body, size=12, color=MUTED)
    add_footer(slide, 8, 15)


def slide_09_principles(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "Our constitution.")
    add_subtitle(slide, "Five principles. We violate none of them. We have a constitution.")

    principles = [
        ("1", "Spec-first development", "Every requirement has an ID; every test references it."),
        ("2", "Clean Architecture", "Domain logic is framework-free, audited by automated tests."),
        ("3", "Test-driven with AI eval gates", "AI outputs scored before they reach users."),
        ("4", "Multi-tenant isolation by construction", "Database-level RLS."),
        ("5", "Observability and cost awareness", "Every request has a correlation ID; cheapest adequate model first."),
    ]
    for i, (num, head, body) in enumerate(principles):
        y = 2.2 + i * 0.95
        add_text_box(slide, 0.7, y, 0.6, 0.7, num, size=36, color=CORAL, bold=True)
        add_text_box(slide, 1.6, y, 4.0, 0.5, head, size=18, color=INDIGO, bold=True)
        add_text_box(slide, 5.6, y, 7.0, 0.7, body, size=14, color=NEAR_BLACK)
    add_footer(slide, 9, 15)


def slide_10_stack(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "Standard. Modern. Mature.")
    add_subtitle(slide, "We are not inventing new technology; we are applying it carefully.")

    rows = [
        ("Backend", "PHP 8.3, Laravel 11 (Clean Architecture)"),
        ("Frontend", "Angular 18 (standalone, signals), NG-ZORRO"),
        ("AI service", "Python 3.12, FastAPI, Pydantic v2"),
        ("Database", "PostgreSQL 16 + pgvector"),
        ("Cache & locks", "Redis 7"),
        ("Object storage", "S3-compatible (MinIO dev / Azure Blob prod)"),
        ("Auth", "OIDC / OAuth2 (Auth0)"),
        ("AI models", "OpenAI, Anthropic, open models (cost-aware router)"),
        ("Container / orch", "Docker + Kubernetes (AKS)"),
        ("IaC", "Terraform"),
        ("CI/CD", "GitHub Actions (lint, tests, SAST, SCA, secrets, eval gates)"),
        ("Tracing", "OpenTelemetry → Jaeger / Azure Monitor"),
    ]
    y = 2.2
    for i, (layer, tech) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else WARM_WHITE
        add_rect(slide, 0.7, y, 12.0, 0.4, bg)
        add_text_box(slide, 0.9, y, 3.0, 0.4, layer, size=13, color=INDIGO, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, 4.0, y, 8.5, 0.4, tech, size=13, color=NEAR_BLACK, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.4
    add_footer(slide, 10, 15)


def slide_11_metrics(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "What we measure.")
    add_subtitle(slide, "These are the four metrics our board will see every month.")

    cells = [
        ("SC-001", "< 5 s", "Discovery p95"),
        ("SC-002", "< 6 s", "Draft p95"),
        ("SC-003", "0.85 / 0.90", "AI relevance / faithfulness"),
        ("SC-004", "99.9%", "API availability (reads)"),
    ]
    for i, (sid, val, label) in enumerate(cells):
        x = 0.7 + i * 3.1
        add_rounded(slide, x, 2.3, 2.9, 2.5, INDIGO)
        add_text_box(slide, x, 2.5, 2.9, 0.4, sid, size=12, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 3.0, 2.9, 1.0, val, size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, x + 0.2, 4.0, 2.5, 0.7, label, size=12, color=WARM_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(
        slide,
        0.7,
        5.4,
        12.0,
        0.4,
        "SC-005 — 0 cross-tenant data leaks (under automated testing)",
        size=14,
        color=NEAR_BLACK,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        0.7,
        5.9,
        12.0,
        0.4,
        "SC-007 — < 10 min to set up a new tenant",
        size=14,
        color=NEAR_BLACK,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 11, 15)


def slide_12_cost(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "AI costs are first-class.")

    tiers = [
        ("Free", "Discovery + 5 boilerplate docs"),
        ("Team", "$99 / mo + AI pass-through\n(most tenants)"),
        ("Pro", "$499 / mo + dedicated model tier\n+ custom funder rules"),
    ]
    for i, (name, desc) in enumerate(tiers):
        x = 0.7 + i * 4.1
        add_rounded(slide, x, 2.3, 3.9, 2.5, WHITE, line=INDIGO)
        add_rect(slide, x, 2.3, 3.9, 0.1, INDIGO)
        add_text_box(slide, x, 2.6, 3.9, 0.5, name, size=22, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.3, 3.4, 3.3, 1.5, desc, size=14, color=NEAR_BLACK, align=PP_ALIGN.CENTER)

    add_text_box(
        slide,
        0.7,
        5.2,
        12.0,
        0.4,
        "Cost-control guarantees",
        size=18,
        color=INDIGO,
        bold=True,
    )
    add_bullets(
        slide,
        0.7,
        5.6,
        12.0,
        1.5,
        [
            "Cost-aware router: cheapest adequate model first.",
            "Per-tenant FinOps dashboard with cost ceilings and alerts.",
            "GPU scale-to-zero when the AI service is idle.",
            "Automatic model fallback on provider outage (no vendor lock-in).",
        ],
        size=14,
    )
    add_footer(slide, 12, 15)


def slide_13_roadmap(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "Where we're going.")

    phases = [
        ("Q3 2026", "Scaffold + Foundation", "Phase 1–2 complete\nPhase 3 (US1 Discovery) in flight"),
        ("Q4 2026", "Org Profile + Library\nProposal Drafting", "Internal alpha\nEnd-to-end P1 MVP"),
        ("Q1 2027", "Budget Narrative\nTracking", "Closed beta with 10 nonprofits"),
        ("Q2 2027", "Funder Tailoring\nReviewer Workflow", "Public GA"),
    ]
    x = 0.7
    box_w = 2.9
    gap = 0.15
    for i, (q, head, body) in enumerate(phases):
        add_rounded(slide, x, 2.3, box_w, 4.0, WHITE, line=INDIGO)
        add_rect(slide, x, 2.3, box_w, 0.1, CORAL)
        add_text_box(slide, x, 2.5, box_w, 0.4, q, size=14, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 3.0, box_w - 0.4, 1.6, head, size=16, color=INDIGO, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 4.7, box_w - 0.4, 1.4, body, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        x += box_w + gap
    add_footer(slide, 13, 15)


def slide_14_ask(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WARM_WHITE)
    add_title(slide, "What we need from you.")

    asks = [
        ("10 design partners", "Closed beta Q1 2027\nYour team, your RFPs, your feedback"),
        ("A reference call", "Introduce us to a foundation program officer or nonprofit ED"),
        ("Co-marketing", "Case study + a quote for the public launch (Q2 2027)"),
    ]
    for i, (head, body) in enumerate(asks):
        y = 2.3 + i * 1.5
        add_rounded(slide, 0.7, y, 8.0, 1.3, WHITE, line=INDIGO)
        add_rect(slide, 0.7, y, 0.15, 1.3, CORAL)
        add_text_box(slide, 1.0, y + 0.15, 7.5, 0.5, head, size=20, color=INDIGO, bold=True)
        add_text_box(slide, 1.0, y + 0.6, 7.5, 0.7, body, size=14, color=NEAR_BLACK)

    # Contact card
    add_rounded(slide, 9.2, 2.3, 3.4, 4.0, INDIGO)
    add_text_box(slide, 9.4, 2.6, 3.0, 0.4, "Contact", size=12, color=CORAL, bold=True)
    add_text_box(slide, 9.4, 3.1, 3.0, 0.5, "[Your Name]", size=20, color=WHITE, bold=True)
    add_text_box(slide, 9.4, 3.7, 3.0, 0.4, "[Your Role]", size=14, color=WARM_WHITE)
    add_text_box(slide, 9.4, 4.3, 3.0, 0.4, "[email@grantgenie.example]", size=12, color=WARM_WHITE)
    add_text_box(slide, 9.4, 4.9, 3.0, 0.4, "[calendar link]", size=12, color=WARM_WHITE)
    add_footer(slide, 14, 15)


def slide_15_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, INDIGO)
    add_rect(slide, 0, 0, 13.333, 7.5, INDIGO)
    add_rect(slide, 0, 0, 5.0, 7.5, INDIGO_LIGHT)

    add_text_box(
        slide,
        0.8,
        2.6,
        11.5,
        1.4,
        "Fund the missions that fund us all.",
        size=48,
        color=WHITE,
        bold=True,
    )
    add_text_box(
        slide,
        0.8,
        4.2,
        11.5,
        0.7,
        "GrantGenie turns weeks of proposal work into hours — with citations your board can trust.",
        size=18,
        color=WARM_WHITE,
    )
    add_text_box(
        slide,
        0.8,
        6.4,
        11.5,
        0.5,
        "Thank you.  Questions?",
        size=18,
        color=CORAL,
        bold=True,
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_why_now(prs)
    slide_04_maya(prs)
    slide_05_solution(prs)
    slide_06_how_it_works(prs)
    slide_07_trust(prs)
    slide_08_ux(prs)
    slide_09_principles(prs)
    slide_10_stack(prs)
    slide_11_metrics(prs)
    slide_12_cost(prs)
    slide_13_roadmap(prs)
    slide_14_ask(prs)
    slide_15_closing(prs)

    out = Path(__file__).resolve().parent / "GrantGenie-Presentation.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
